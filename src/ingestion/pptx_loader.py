import base64
import io
from pathlib import Path

from pptx import Presentation
from PIL import Image

from config import IMAGE_MAX_SIZE
from logger import get_logger
from .base import DocumentChunk

log = get_logger("ingestion.pptx")


def load_pptx(path: str) -> list[DocumentChunk]:
    chunks: list[DocumentChunk] = []
    source = Path(path).name
    log.info("Opening PPTX: %s", source)

    prs = Presentation(path)
    total_slides = len(prs.slides)
    log.info("Presentation has %d slide(s)", total_slides)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text, slide_images = [], 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)

            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    img.thumbnail(IMAGE_MAX_SIZE)
                    buf = io.BytesIO()
                    img.save(buf, format="JPEG", quality=85)
                    b64 = base64.b64encode(buf.getvalue()).decode()
                    chunks.append(DocumentChunk(
                        content="",
                        chunk_type="image",
                        source=source,
                        page=slide_num,
                        metadata={"file_type": "pptx", "slide": slide_num},
                        image_b64=b64,
                    ))
                    slide_images += 1
                except Exception as e:
                    log.warning("Slide %d: failed to extract image: %s", slide_num, e)

        if slide_text:
            chunks.append(DocumentChunk(
                content="\n".join(slide_text),
                chunk_type="text",
                source=source,
                page=slide_num,
                metadata={"file_type": "pptx", "slide": slide_num},
            ))

        log.info("Slide %d/%d: text_lines=%d  images=%d",
                 slide_num, total_slides, len(slide_text), slide_images)

    log.info("PPTX load complete | total_chunks=%d", len(chunks))
    return chunks
